"""Tạo 2 slide: Encoder Pipeline và Decoder Pipeline."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Màu sắc
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)
ACCENT_ORANGE = RGBColor(0xFF, 0x57, 0x22)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xBB, 0xBB, 0xBB)
DARK_CARD = RGBColor(0x25, 0x25, 0x40)
HIGHLIGHT = RGBColor(0xFF, 0xD5, 0x4F)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, left, top, w, h, fill_color, text, font_size=14,
            font_color=WHITE, bold=False, border_color=None):
    """Thêm hộp với text."""
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(
        1, Inches(left), Inches(top), Inches(w), Inches(h))  # 1 = Rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()

    # Bo góc
    spPr = shape._element.spPr
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is not None:
        prstGeom.set('prst', 'roundRect')

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    # Vertical center
    tf.word_wrap = True
    from pptx.enum.text import MSO_ANCHOR
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT_BLUE, width=Pt(3)):
    """Thêm mũi tên."""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = width
    # Arrow head
    from pptx.oxml.ns import qn
    ln = connector._element.find(qn('a:ln'))
    if ln is None:
        from lxml import etree
        ln = etree.SubElement(connector._element.spPr, qn('a:ln'))
    tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tail)


def add_text(slide, left, top, w, h, text, size=12, color=LIGHT_GRAY, bold=False, align=PP_ALIGN.LEFT):
    """Thêm text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


# ======================================================================
# SLIDE 1: ENCODER PIPELINE
# ======================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
set_slide_bg(slide1, DARK_BG)

# Title
add_text(slide1, 0.5, 0.2, 12, 0.7, "ENCODER PIPELINE", 32, WHITE, True, PP_ALIGN.CENTER)
add_text(slide1, 0.5, 0.7, 12, 0.5, "Cross-Layer Adaptive QP Video Encoding", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ---- Row 1: Main Pipeline Flow ----
y_main = 1.6
box_h = 0.9

# Box 1: Video Input
add_box(slide1, 0.3, y_main, 1.6, box_h, DARK_CARD, "📹 Video YUV\nForeman CIF\n352×288", 11, WHITE, False, ACCENT_BLUE)

# Arrow 1→2
add_text(slide1, 1.95, y_main + 0.25, 0.5, 0.4, "→", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Box 2: BCD+SCA Optimizer
add_box(slide1, 2.4, y_main, 2.2, box_h, RGBColor(0x0D, 0x47, 0xA1), "🔧 BCD+SCA\nOptimizer\n(8 iterations)", 12, WHITE, True, ACCENT_BLUE)

# Arrow 2→3
add_text(slide1, 4.65, y_main + 0.25, 0.5, 0.4, "→", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Box 3: Output = P(t), q(t), R(t)
add_box(slide1, 5.1, y_main, 1.8, box_h, DARK_CARD, "P(t), q(t)\n→ R(t)\nChannel Rate", 11, HIGHLIGHT, True, ACCENT_ORANGE)

# Arrow 3→4
add_text(slide1, 6.95, y_main + 0.25, 0.5, 0.4, "→", 24, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# Box 4: QP Mapping
add_box(slide1, 7.4, y_main, 2.0, box_h, RGBColor(0xE6, 0x51, 0x00), "📐 QP Mapping\nQP = 32-(R-2)×11\nRange: 10~32", 11, WHITE, True, ACCENT_ORANGE)

# Arrow 4→5
add_text(slide1, 9.45, y_main + 0.25, 0.5, 0.4, "→", 24, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# Box 5: H.264 Encoder
add_box(slide1, 9.9, y_main, 1.8, box_h, RGBColor(0x1B, 0x5E, 0x20), "🎬 H.264\nFFmpeg\nEncode", 12, WHITE, True, ACCENT_GREEN)

# Arrow 5→6
add_text(slide1, 11.75, y_main + 0.25, 0.5, 0.4, "→", 24, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# Box 6: Bitstream
add_box(slide1, 12.1, y_main, 1.0, box_h, DARK_CARD, "💾\n.264", 14, ACCENT_GREEN, True, ACCENT_GREEN)

# ---- Row 2: Details ----
y_detail = 3.0

# Detail: BCD+SCA
add_box(slide1, 0.3, y_detail, 4.0, 2.8, DARK_CARD, "", 10, WHITE, False, ACCENT_BLUE)
add_text(slide1, 0.5, y_detail + 0.1, 3.6, 0.4, "⚙️ BCD+SCA Chi Tiết", 14, ACCENT_BLUE, True)
add_text(slide1, 0.5, y_detail + 0.5, 3.6, 2.2,
         "Bước 1: Cố định q(t) → Tối ưu P(t)\n"
         "   • Giải closed-form (waterfilling)\n"
         "   • Ràng buộc: P ≤ Pmax, I ≤ Ith\n\n"
         "Bước 2: Cố định P(t) → Tối ưu q(t)\n"
         "   • SCA xấp xỉ Taylor bậc 1\n"
         "   • SLSQP solver\n"
         "   • Ràng buộc: tốc độ bay ≤ Vmax\n\n"
         "Lặp 8 vòng → QoE hội tụ",
         10, LIGHT_GRAY)

# Detail: Adaptive QP vs SVC
add_box(slide1, 4.6, y_detail, 4.0, 2.8, DARK_CARD, "", 10, WHITE, False, ACCENT_ORANGE)
add_text(slide1, 4.8, y_detail + 0.1, 3.6, 0.4, "🔄 Adaptive QP vs SVC", 14, ACCENT_ORANGE, True)
add_text(slide1, 4.8, y_detail + 0.5, 3.6, 2.2,
         "Adaptive QP (đề xuất):\n"
         "  • QP thay đổi LIÊN TỤC theo R(t)\n"
         "  • Tận dụng 100% băng thông\n"
         "  • Encoder quyết định chất lượng\n\n"
         "SVC (truyền thống):\n"
         "  • QP CỐ ĐỊNH: BL=35, EL=15\n"
         "  • Chất lượng rời rạc (3 mức)\n"
         "  • Decoder quyết định chất lượng\n"
         "  • Lãng phí khi R gần ngưỡng",
         10, LIGHT_GRAY)

# Detail: Formulas
add_box(slide1, 8.9, y_detail, 4.2, 2.8, DARK_CARD, "", 10, WHITE, False, ACCENT_GREEN)
add_text(slide1, 9.1, y_detail + 0.1, 3.8, 0.4, "📊 Công Thức Chính", 14, ACCENT_GREEN, True)
add_text(slide1, 9.1, y_detail + 0.5, 3.8, 2.2,
         "Channel Rate:\n"
         "  R(t) = log₂(1 + P(t)·h(q(t)) / σ²)\n\n"
         "Channel Gain:\n"
         "  h(q) = β₀ / (‖q - wBS‖² + H²)\n\n"
         "QoE Function (Log):\n"
         "  QoE(R) = log₂(1 + R)\n\n"
         "Objective:\n"
         "  max Σ QoE(R(t))  s.t. constraints",
         10, LIGHT_GRAY)

# Footer
add_text(slide1, 0.3, 6.5, 12, 0.5,
         "💡 Key Insight: Kênh truyền (R) trực tiếp điều khiển chất lượng nén (QP) → Cross-layer optimization",
         13, HIGHLIGHT, True, PP_ALIGN.CENTER)


# ======================================================================
# SLIDE 2: DECODER PIPELINE
# ======================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2, DARK_BG)

# Title
add_text(slide2, 0.5, 0.2, 12, 0.7, "DECODER & EVALUATION PIPELINE", 32, WHITE, True, PP_ALIGN.CENTER)
add_text(slide2, 0.5, 0.7, 12, 0.5, "H.264 Decoding → PSNR Quality Assessment", 16, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# ---- Row 1: Decode Pipeline ----
y_main = 1.5
box_h = 0.85

# Step 1
add_box(slide2, 0.2, y_main, 1.4, box_h, DARK_CARD, "💾 Bitstream\n.264 file", 11, ACCENT_BLUE, True, ACCENT_BLUE)
add_text(slide2, 1.6, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Step 2
add_box(slide2, 2.0, y_main, 1.6, box_h, RGBColor(0x0D, 0x47, 0xA1), "① Parse NAL\nSPS/PPS\nSlice Header", 10, WHITE, False, ACCENT_BLUE)
add_text(slide2, 3.6, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Step 3
add_box(slide2, 4.0, y_main, 1.6, box_h, RGBColor(0x0D, 0x47, 0xA1), "② Entropy\nDecode\nCABAC/CAVLC", 10, WHITE, False, ACCENT_BLUE)
add_text(slide2, 5.6, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Step 4
add_box(slide2, 6.0, y_main, 1.6, box_h, RGBColor(0xE6, 0x51, 0x00), "③ Inverse\nQuantize\ncoeff × Qstep", 10, WHITE, True, ACCENT_ORANGE)
add_text(slide2, 7.6, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_ORANGE, True, PP_ALIGN.CENTER)

# Step 5
add_box(slide2, 8.0, y_main, 1.6, box_h, RGBColor(0xE6, 0x51, 0x00), "④ IDCT 4×4\nInteger\nTransform", 10, WHITE, False, ACCENT_ORANGE)
add_text(slide2, 9.6, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# Step 6
add_box(slide2, 10.0, y_main, 1.5, box_h, RGBColor(0x1B, 0x5E, 0x20), "⑤ Prediction\n+ Residual", 10, WHITE, False, ACCENT_GREEN)
add_text(slide2, 11.5, y_main + 0.2, 0.4, 0.4, "→", 22, ACCENT_GREEN, True, PP_ALIGN.CENTER)

# Step 7
add_box(slide2, 11.9, y_main, 1.2, box_h, RGBColor(0x1B, 0x5E, 0x20), "⑥ Deblock\nFilter", 10, WHITE, False, ACCENT_GREEN)

# QP annotation
add_box(slide2, 5.5, y_main - 0.7, 2.6, 0.5, RGBColor(0x4A, 0x14, 0x8C),
        "QP đọc từ bitstream — decoder tự biết!", 9, HIGHLIGHT, True, ACCENT_PURPLE)

# ---- Row 2: SVC Decode Logic ----
y_svc = 3.0

add_box(slide2, 0.2, y_svc, 6.2, 2.6, DARK_CARD, "", 10, WHITE, False, ACCENT_PURPLE)
add_text(slide2, 0.4, y_svc + 0.1, 5.8, 0.4, "📡 SVC Decode Logic (3 mức chất lượng)", 14, ACCENT_PURPLE, True)

# 3 levels
y_lv = y_svc + 0.6
add_box(slide2, 0.4, y_lv, 5.8, 0.55, RGBColor(0x1B, 0x5E, 0x20),
        "R(t) ≥ 3.0 bps/Hz  →  BL + EL  →  Chất lượng CAO (QP_eff ≈ 15)", 11, WHITE, True)

add_box(slide2, 0.4, y_lv + 0.65, 5.8, 0.55, RGBColor(0xE6, 0x51, 0x00),
        "R(t) ≥ 2.0 bps/Hz  →  BL only  →  MỜ nhưng xem được (QP = 35)", 11, WHITE, True)

add_box(slide2, 0.4, y_lv + 1.3, 5.8, 0.55, RGBColor(0xB7, 0x1C, 0x1C),
        "R(t) < 2.0 bps/Hz  →  Freeze frame  →  ĐỨNG HÌNH ❌", 11, WHITE, True)

# ---- Row 2 right: PSNR Evaluation ----
add_box(slide2, 6.7, y_svc, 6.2, 2.6, DARK_CARD, "", 10, WHITE, False, ACCENT_GREEN)
add_text(slide2, 6.9, y_svc + 0.1, 5.8, 0.4, "📊 Đánh Giá Chất Lượng — PSNR", 14, ACCENT_GREEN, True)
add_text(slide2, 6.9, y_svc + 0.55, 5.8, 2.0,
         "PSNR = 10 × log₁₀(255² / MSE)\n\n"
         "• So sánh từng frame: decoded vs original YUV\n"
         "• Full-Reference metric (cần video gốc)\n"
         "• Đơn vị: dB (cao hơn = tốt hơn)\n\n"
         "Kết quả:\n"
         "  Adaptive QP:  33.6 ~ 41.2 dB  ✅ Tốt nhất\n"
         "  SVC:          27.5 ~ 39.0 dB\n"
         "  Straight:     27.2 ~ 30.0 dB\n"
         "  Circle:       26.9 ~ 31.2 dB",
         10, LIGHT_GRAY)

# ---- Row 3: Key Results ----
y_res = 5.9
add_box(slide2, 0.2, y_res, 4.2, 1.1, DARK_CARD, "", 10, WHITE, False, ACCENT_BLUE)
add_text(slide2, 0.4, y_res + 0.05, 3.8, 0.35, "🏆 Kết quả chính", 13, ACCENT_BLUE, True)
add_text(slide2, 0.4, y_res + 0.4, 3.8, 0.6,
         "Adaptive QP vượt Straight: +6 ~ +11 dB\n"
         "Adaptive QP vượt SVC: +2 ~ +6 dB\n"
         "Adaptive QP luôn CAO NHẤT mọi case",
         10, HIGHLIGHT, True)

add_box(slide2, 4.6, y_res, 4.2, 1.1, DARK_CARD, "", 10, WHITE, False, ACCENT_ORANGE)
add_text(slide2, 4.8, y_res + 0.05, 3.8, 0.35, "🔑 Tại sao Adaptive QP tốt hơn?", 13, ACCENT_ORANGE, True)
add_text(slide2, 4.8, y_res + 0.4, 3.8, 0.6,
         "• QP thay đổi → tận dụng 100% R(t)\n"
         "• SVC lãng phí khi R gần ngưỡng\n"
         "• SVC freeze khi R < 2.0 → PSNR rất thấp",
         10, LIGHT_GRAY)

add_box(slide2, 9.0, y_res, 4.0, 1.1, DARK_CARD, "", 10, WHITE, False, ACCENT_GREEN)
add_text(slide2, 9.2, y_res + 0.05, 3.6, 0.35, "📌 Lưu ý kỹ thuật", 13, ACCENT_GREEN, True)
add_text(slide2, 9.2, y_res + 0.4, 3.6, 0.6,
         "• Decoder KHÔNG cần biết QP\n"
         "• QP được nhúng trong bitstream\n"
         "• PSNR là Full-Reference metric",
         10, LIGHT_GRAY)

# Save
out = Path(__file__).parent.parent / "results" / "Encoder_Decoder_Slides.pptx"
prs.save(str(out))
print(f"Saved: {out}")
