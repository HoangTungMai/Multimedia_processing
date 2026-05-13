"""
Tạo slide PowerPoint cho đề tài UAV CR-NOMA Video Transmission
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === Color scheme ===
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD = RGBColor(0x25, 0x25, 0x40)
ACCENT_BLUE = RGBColor(0x21, 0x96, 0xF3)
ACCENT_ORANGE = RGBColor(0xFF, 0x57, 0x22)
ACCENT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_PURPLE = RGBColor(0x9C, 0x27, 0xB0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0xB0, 0xB0, 0xB0)

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def add_para(tf, text, size=16, color=WHITE, bold=False, space_before=6):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(space_before)
    return p

def add_card(slide, left, top, width, height, color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_image_safe(slide, path, left, top, width, height):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))
    else:
        add_text(slide, left, top, width, height, f"[Missing: {os.path.basename(path)}]", 12, GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
# Accent line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT_BLUE; shape.line.fill.background()

add_text(slide, 1, 1.5, 11, 1.2, "TỐI ƯU TRUYỀN VIDEO UAV", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 2.7, 11, 1, "TRÊN MẠNG CR-NOMA", 40, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 4.2, 11, 0.8, "Cross-layer Adaptive Video Encoding với BCD+SCA Trajectory Optimization", 20, GRAY, align=PP_ALIGN.CENTER)
add_text(slide, 1, 5.5, 11, 0.6, "Môn: Xử lý ảnh và Truyền thông đa phương tiện", 18, WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 1, 6.2, 11, 0.6, "GVHD: ________________    |    SV: ________________", 16, GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: Bài toán
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 5, 0.7, "01  BÀI TOÁN ĐẶT RA", 28, ACCENT_BLUE, True)

add_card(slide, 0.5, 1.2, 6, 5.5)
tf = add_text(slide, 0.8, 1.4, 5.5, 5, "Bối cảnh", 22, ACCENT_BLUE, True)
add_para(tf, "• UAV bay từ A → B, vừa quay video vừa truyền về BS", 16, WHITE)
add_para(tf, "• Trên mặt đất có PU (Primary User) cùng tần số", 16, WHITE)
add_para(tf, "• UAV phải chia sẻ phổ mà không gây nhiễu cho PU", 16, WHITE)
add_para(tf, "", 10)
add_para(tf, "Mục tiêu", 22, ACCENT_ORANGE, True)
add_para(tf, "Tối đa chất lượng video (PSNR) với ràng buộc:", 16, WHITE)
add_para(tf, "  ✦ Tốc độ bay: Vmax = 15 m/s", 15, GRAY)
add_para(tf, "  ✦ Công suất phát: Pmax = 0.5 W", 15, GRAY)
add_para(tf, "  ✦ Nhiễu CR: I_th = 5×10⁻¹¹ W", 15, GRAY)
add_para(tf, "  ✦ QoS cho PU: R_PU,min = 0.5 bps/Hz", 15, GRAY)

add_card(slide, 7, 1.2, 5.8, 5.5)
add_text(slide, 7.3, 1.4, 5.2, 0.5, "Mô hình hệ thống", 20, ACCENT_BLUE, True)
add_image_safe(slide, "results/trajectories_all_cases.png", 7.2, 2.0, 5.4, 4.2)

# ============================================================
# SLIDE 3: Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 8, 0.7, "02  QUY TRÌNH THỰC HIỆN", 28, ACCENT_BLUE, True)

# Step 1
add_card(slide, 0.5, 1.3, 3.8, 5.5)
add_text(slide, 0.8, 1.5, 3.3, 0.5, "BƯỚC 1: Tối ưu quỹ đạo", 18, ACCENT_ORANGE, True)
tf = add_text(slide, 0.8, 2.1, 3.3, 4, "BCD+SCA Optimizer", 16, WHITE, True)
add_para(tf, "• Lặp 8 vòng BCD", 14, GRAY)
add_para(tf, "• ① Tối ưu công suất P(t)", 14, GRAY)
add_para(tf, "• ② Tối ưu quỹ đạo q(t) bằng SCA", 14, GRAY)
add_para(tf, "• Xấp xỉ Taylor bậc 1 + SLSQP", 14, GRAY)
add_para(tf, "", 8)
add_para(tf, "→ Đầu ra: rates[30], trajectory, power", 14, ACCENT_GREEN)

# Step 2
add_card(slide, 4.6, 1.3, 3.8, 5.5)
add_text(slide, 4.9, 1.5, 3.3, 0.5, "BƯỚC 2: Mã hóa video", 18, ACCENT_ORANGE, True)
tf = add_text(slide, 4.9, 2.1, 3.3, 4, "H.264 Encoder (FFmpeg)", 16, WHITE, True)
add_para(tf, "• Video: foreman_cif.yuv", 14, GRAY)
add_para(tf, "• 300 frames → 30 GOP × 10 frames", 14, GRAY)
add_para(tf, "• Adaptive QP: liên tục theo R(t)", 14, GRAY)
add_para(tf, "• SVC: BL (QP=35) + EL (QP=15)", 14, GRAY)
add_para(tf, "", 8)
add_para(tf, "→ Đầu ra: .264 bitstream", 14, ACCENT_GREEN)

# Step 3
add_card(slide, 8.7, 1.3, 4.1, 5.5)
add_text(slide, 9.0, 1.5, 3.6, 0.5, "BƯỚC 3: Đánh giá PSNR", 18, ACCENT_ORANGE, True)
tf = add_text(slide, 9.0, 2.1, 3.6, 4, "Full-Reference Evaluation", 16, WHITE, True)
add_para(tf, "• Decode .264 → frames", 14, GRAY)
add_para(tf, "• So pixel-by-pixel với gốc", 14, GRAY)
add_para(tf, "• PSNR = 20·log₁₀(255/√MSE)", 14, GRAY)
add_para(tf, "• 5 phương pháp × 5 cấu hình", 14, GRAY)
add_para(tf, "", 8)
add_para(tf, "→ Đầu ra: biểu đồ so sánh", 14, ACCENT_GREEN)

# ============================================================
# SLIDE 4: Cross-layer Design
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "03  CROSS-LAYER DESIGN: Adaptive QP vs SVC", 28, ACCENT_BLUE, True)

# Left: Adaptive QP
add_card(slide, 0.5, 1.3, 5.8, 5.5)
add_text(slide, 0.8, 1.5, 5.3, 0.5, "✦ ADAPTIVE QP (Đề xuất)", 20, ACCENT_GREEN, True)
tf = add_text(slide, 0.8, 2.2, 5.3, 4, "Ánh xạ Cross-layer: QP = 32 − (R−2.0) × 11", 14, WHITE, True)
add_para(tf, "", 6)
add_para(tf, "R(t) = 2.0 bps/Hz  →  QP = 32  →  Mờ", 14, GRAY)
add_para(tf, "R(t) = 2.5 bps/Hz  →  QP = 27  →  Trung bình", 14, GRAY)
add_para(tf, "R(t) = 3.0 bps/Hz  →  QP = 21  →  Khá tốt", 14, GRAY)
add_para(tf, "R(t) = 3.5 bps/Hz  →  QP = 16  →  Rất nét", 14, GRAY)
add_para(tf, "R(t) = 4.0 bps/Hz  →  QP = 10  →  Gần lossless", 14, GRAY)
add_para(tf, "", 8)
add_para(tf, "✅ Ưu điểm: Tận dụng tối đa mỗi bps/Hz", 15, ACCENT_GREEN)
add_para(tf, "✅ QP thay đổi liên tục → chất lượng liên tục", 15, ACCENT_GREEN)
add_para(tf, "❌ Nhược: Không chống mất gói, cần CSI chính xác", 15, ACCENT_ORANGE)

# Right: SVC
add_card(slide, 6.7, 1.3, 6.1, 5.5)
add_text(slide, 7.0, 1.5, 5.5, 0.5, "✦ SVC (Truyền thống)", 20, ACCENT_PURPLE, True)
tf = add_text(slide, 7.0, 2.2, 5.5, 4, "2 lớp cố định: BL (QP=35) + EL (QP=15)", 14, WHITE, True)
add_para(tf, "", 6)
add_para(tf, "R(t) ≥ 3.0 → Dùng EL (QP=15, nét)", 14, GRAY)
add_para(tf, "R(t) ≥ 2.0 → Dùng BL (QP=35, mờ)", 14, GRAY)
add_para(tf, "R(t) < 2.0 → Freeze frame trước", 14, GRAY)
add_para(tf, "", 8)
add_para(tf, "Ví dụ: Khi R(t) = 2.8 bps/Hz", 16, WHITE, True)
add_para(tf, "• Adaptive QP: QP=23, PSNR ~38 dB", 14, ACCENT_GREEN)
add_para(tf, "• SVC: chỉ BL (QP=35), PSNR ~28 dB", 14, ACCENT_PURPLE)
add_para(tf, "→ Lãng phí 0.8 bps/Hz!", 14, ACCENT_ORANGE)
add_para(tf, "", 8)
add_para(tf, "✅ Ưu: Chống mất gói, hỗ trợ multicast", 15, ACCENT_GREEN)
add_para(tf, "❌ Nhược: Chất lượng rời rạc (bậc thang)", 15, ACCENT_ORANGE)

# ============================================================
# SLIDE 5: Hàm QoE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "04  HÀM MỤC TIÊU QoE", 28, ACCENT_BLUE, True)

add_card(slide, 0.5, 1.3, 6, 2.5)
add_text(slide, 0.8, 1.5, 5.5, 0.5, "Sigmoid QoE (cho SVC)", 20, ACCENT_BLUE, True)
tf = add_text(slide, 0.8, 2.1, 5.5, 1.5, "QoE = 0.6/(1+e^(-5(R-2.0))) + 0.4/(1+e^(-5(R-3.0)))", 16, WHITE)
add_para(tf, "Hai ngưỡng: R_BL=2.0 và R_BL+R_EL=3.0", 14, GRAY)
add_para(tf, "Bão hòa khi vượt ngưỡng → không khuyến khích tăng thêm", 14, GRAY)

add_card(slide, 0.5, 4.2, 6, 2.5)
add_text(slide, 0.8, 4.4, 5.5, 0.5, "Logarithmic QoE (cho Adaptive QP)", 20, ACCENT_ORANGE, True)
tf = add_text(slide, 0.8, 5.0, 5.5, 1.5, "QoE = log₂(1 + R)", 16, WHITE)
add_para(tf, "Luôn tăng khi Rate tăng, không bão hòa", 14, GRAY)
add_para(tf, "Mỗi bps/Hz thêm đều có giá trị", 14, GRAY)

add_card(slide, 7, 1.3, 5.8, 5.4)
add_text(slide, 7.3, 1.5, 5.2, 0.5, "BCD Convergence", 20, ACCENT_GREEN, True)
add_image_safe(slide, "results/bcd_convergence.png", 7.1, 2.1, 5.6, 4.3)

# ============================================================
# SLIDE 6: Kết quả PSNR
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "05  KẾT QUẢ PSNR", 28, ACCENT_BLUE, True)

add_image_safe(slide, "results/psnr_summary_bar.png", 0.3, 1.2, 6.5, 5.5)

add_card(slide, 7, 1.2, 5.8, 5.5)
add_text(slide, 7.3, 1.4, 5.2, 0.5, "Bảng kết quả (dB)", 20, ACCENT_BLUE, True)
data = [
    ("Case", "Sigmoid", "LogQoE", "SVC", "Straight"),
    ("Case1", "37.87", "37.97", "35.28", "28.76"),
    ("Case2", "37.90", "37.97", "35.57", "28.76"),
    ("Case3", "33.59", "33.59", "27.54", "27.16"),
    ("Case4", "41.09", "41.16", "38.95", "29.97"),
    ("Case5", "38.97", "39.04", "33.44", "28.76"),
]
tf = add_text(slide, 7.3, 2.0, 5.2, 4.5, "", 12, WHITE)
for i, row in enumerate(data):
    line = f"{'  '.join(f'{c:<10}' for c in row)}"
    color = ACCENT_BLUE if i == 0 else WHITE
    bold = (i == 0)
    add_para(tf, line, 12, color, bold, space_before=4)

add_para(tf, "", 10)
add_para(tf, "✦ BCD+SCA: +6 đến +11 dB vs Straight", 14, ACCENT_GREEN, True)
add_para(tf, "✦ Adaptive QP ≥ SVC trong mọi case", 14, ACCENT_GREEN, True)
add_para(tf, "✦ Log QoE ≈ Sigmoid (sau fix R_EL)", 14, ACCENT_ORANGE, True)

# ============================================================
# SLIDE 7: So sánh tổng hợp
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "06  SO SÁNH TỔNG HỢP", 28, ACCENT_BLUE, True)
add_image_safe(slide, "results/comprehensive_comparison.png", 0.3, 1.0, 12.7, 6.2)

# ============================================================
# SLIDE 8: SVC Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "07  PHÂN TÍCH SVC vs ADAPTIVE QP", 28, ACCENT_BLUE, True)
add_image_safe(slide, "results/svc_optimal_analysis.png", 0.3, 1.0, 6.3, 5.5)
add_image_safe(slide, "results/svc_vs_adaptive_detail.png", 6.8, 1.0, 6.2, 5.5)

# ============================================================
# SLIDE 9: Kết luận
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_text(slide, 0.5, 0.3, 10, 0.7, "08  KẾT LUẬN", 28, ACCENT_BLUE, True)

add_card(slide, 0.5, 1.3, 6, 2.2)
add_text(slide, 0.8, 1.4, 5.5, 0.4, "① Tối ưu quỹ đạo là yếu tố then chốt", 18, ACCENT_GREEN, True)
tf = add_text(slide, 0.8, 1.9, 5.5, 1.3, "BCD+SCA mang lại +6 đến +11 dB so với baseline.", 15, WHITE)
add_para(tf, "Quỹ đạo thông minh né PU, bám BS → Rate cao → Video nét.", 14, GRAY)

add_card(slide, 0.5, 3.8, 6, 2.2)
add_text(slide, 0.8, 3.9, 5.5, 0.4, "② Cross-layer Adaptive QP vượt trội SVC", 18, ACCENT_ORANGE, True)
tf = add_text(slide, 0.8, 4.4, 5.5, 1.3, "Adaptive QP tận dụng liên tục từng bps/Hz.", 15, WHITE)
add_para(tf, "SVC lãng phí băng thông ở vùng giữa các ngưỡng.", 14, GRAY)
add_para(tf, "Kênh xấu (Case 3): Adaptive QP vượt SVC +6.05 dB.", 14, GRAY)

add_card(slide, 7, 1.3, 5.8, 2.2)
add_text(slide, 7.3, 1.4, 5.2, 0.4, "③ Hàm QoE ảnh hưởng đến quỹ đạo", 18, ACCENT_PURPLE, True)
tf = add_text(slide, 7.3, 1.9, 5.2, 1.3, "Sigmoid QoE → quỹ đạo an toàn, cân bằng nhiệm vụ.", 15, WHITE)
add_para(tf, "Log QoE → quỹ đạo bám sát BS, tối đa video.", 14, GRAY)

add_card(slide, 7, 3.8, 5.8, 2.2)
add_text(slide, 7.3, 3.9, 5.2, 0.4, "④ Hướng phát triển", 18, WHITE, True)
tf = add_text(slide, 7.3, 4.4, 5.2, 1.3, "• Thêm SSIM/VMAF cho đánh giá cảm nhận", 14, GRAY)
add_para(tf, "• Mô phỏng kênh Rician fading", 14, GRAY)
add_para(tf, "• Mở rộng multicast (1 UAV → nhiều BS)", 14, GRAY)

# ============================================================
# SLIDE 10: Thank you
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT_BLUE; shape.line.fill.background()

add_text(slide, 1, 2.5, 11, 1, "CẢM ƠN THẦY/CÔ ĐÃ LẮNG NGHE", 40, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, 1, 3.8, 11, 0.8, "Q & A", 60, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Save
out_path = "results/UAV_CR_NOMA_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
